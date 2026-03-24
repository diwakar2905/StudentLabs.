from celery_app import celery_app
from database import SessionLocal
from models import Export, Assignment, Presentation
import json
import os
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

@celery_app.task(bind=True, max_retries=3)
def export_assignment_pdf(self, project_id: int, assignment_id: int):
    """Generate actual PDF file for assignment using ReportLab"""
    db = SessionLocal()
    try:
        # Fetch assignment from database
        assignment = db.query(Assignment).filter(Assignment.id == assignment_id).first()
        if not assignment:
            return {"error": "Assignment not found", "assignment_id": assignment_id}
        
        # Create output directory
        os.makedirs("generated", exist_ok=True)
        file_path = f"generated/assignment_{project_id}_{assignment_id}.pdf"
        
        # Create PDF using ReportLab
        c = canvas.Canvas(file_path, pagesize=A4)
        width, height = A4
        
        # Set up some constants
        margin = 40
        y_position = height - margin
        line_height = 14
        
        # Add title
        c.setFont("Helvetica-Bold", 16)
        title = assignment.title or f"Assignment for Project {project_id}"
        c.drawString(margin, y_position, title)
        y_position -= line_height * 2
        
        # Add content
        c.setFont("Helvetica", 10)
        
        # Split content into lines
        content_lines = assignment.content.split("\n")
        
        for line in content_lines:
            # Check if we need a new page
            if y_position < margin:
                c.showPage()
                y_position = height - margin
                c.setFont("Helvetica", 10)
            
            # Handle markdown headers
            if line.startswith("# "):
                c.setFont("Helvetica-Bold", 14)
                c.drawString(margin, y_position, line.replace("# ", ""))
                c.setFont("Helvetica", 10)
            elif line.startswith("## "):
                c.setFont("Helvetica-Bold", 12)
                c.drawString(margin + 10, y_position, line.replace("## ", ""))
                c.setFont("Helvetica", 10)
            elif line.startswith("### "):
                c.setFont("Helvetica-BoldOblique", 11)
                c.drawString(margin + 20, y_position, line.replace("### ", ""))
                c.setFont("Helvetica", 10)
            elif line.strip():  # Regular text
                # Wrap text if too long
                if len(line) > 100:
                    # Simple text wrapping
                    words = line.split()
                    current_line = ""
                    for word in words:
                        if len(current_line) + len(word) + 1 < 100:
                            current_line += word + " "
                        else:
                            if current_line:
                                c.drawString(margin, y_position, current_line)
                                y_position -= line_height
                            current_line = word + " "
                    if current_line:
                        c.drawString(margin, y_position, current_line)
                else:
                    c.drawString(margin, y_position, line)
            
            y_position -= line_height
        
        # Save PDF
        c.save()
        
        # Create export record in database
        export = Export(
            project_id=project_id,
            file_type="pdf",
            file_path=file_path,
            file_url=f"http://localhost:8000/generated/assignment_{project_id}_{assignment_id}.pdf"
        )
        db.add(export)
        db.commit()
        
        return {
            "status": "completed",
            "project_id": project_id,
            "assignment_id": assignment_id,
            "file_type": "pdf",
            "file_path": file_path,
            "message": f"PDF generated: {file_path}"
        }
    
    except Exception as exc:
        db.rollback()
        return {"error": str(exc), "assignment_id": assignment_id}
    
    finally:
        db.close()

@celery_app.task(bind=True, max_retries=3)
def export_presentation_pptx(self, project_id: int, presentation_id: int):
    """Generate actual PPTX file for presentation using python-pptx"""
    db = SessionLocal()
    try:
        # Fetch presentation from database
        presentation = db.query(Presentation).filter(Presentation.id == presentation_id).first()
        if not presentation:
            return {"error": "Presentation not found", "presentation_id": presentation_id}
        
        # Parse slides JSON
        try:
            slides_data = json.loads(presentation.slides_json) if isinstance(presentation.slides_json, str) else presentation.slides_json
        except (json.JSONDecodeError, TypeError):
            slides_data = presentation.slides_json if isinstance(presentation.slides_json, list) else []
        
        if not slides_data:
            return {"error": "No slides data found", "presentation_id": presentation_id}
        
        # Create output directory
        os.makedirs("generated", exist_ok=True)
        file_path = f"generated/presentation_{project_id}_{presentation_id}.pptx"
        
        # Create PowerPoint presentation
        prs = PptxPresentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Process each slide
        for slide_data in slides_data:
            # Choose layout based on slide type
            layout_type = slide_data.get("layout", "bullet")
            
            if layout_type == "title":
                # Title slide layout
                slide_layout = prs.slide_layouts[0]
            else:
                # Title and content layout
                slide_layout = prs.slide_layouts[1]
            
            # Add slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if len(slide.placeholders) > 0:
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get("title", "")
                title_shape.text_frame.paragraphs[0].font.size = Pt(44)
                title_shape.text_frame.paragraphs[0].font.bold = True
            
            # Set content
            if len(slide.placeholders) > 1:
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                content = slide_data.get("content", [])
                
                if isinstance(content, list):
                    # Bullet points
                    for idx, bullet in enumerate(content):
                        if idx == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = str(bullet)
                        p.level = 0
                        p.font.size = Pt(24)
                
                elif isinstance(content, str):
                    # Single text content
                    text_frame.text = content
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(24)
            
            # Add speaker notes if available
            speaker_notes = slide_data.get("speaker_notes", "")
            if speaker_notes:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                notes_frame.text = speaker_notes
        
        # Save PowerPoint file
        prs.save(file_path)
        
        # Create export record in database
        export = Export(
            project_id=project_id,
            file_type="pptx",
            file_path=file_path,
            file_url=f"http://localhost:8000/generated/presentation_{project_id}_{presentation_id}.pptx"
        )
        db.add(export)
        db.commit()
        
        return {
            "status": "completed",
            "project_id": project_id,
            "presentation_id": presentation_id,
            "file_type": "pptx",
            "file_path": file_path,
            "slides_count": len(slides_data),
            "message": f"PowerPoint generated: {file_path}"
        }
    
    except Exception as exc:
        db.rollback()
        return {"error": str(exc), "presentation_id": presentation_id}
    
    finally:
        db.close()
