"""
PPT Builder Module — Convert Assignments to PowerPoint Slides

This module takes a completed assignment and automatically converts it into
professional presentation slides by:

1. Extracting markdown sections from assignment text
2. Summarizing key points from each section
3. Converting to bullet points
4. Creating slide structure
5. Generating speaker notes

Architecture:
Assignment (Markdown) → Sections → Summaries → Bullets → Slide JSON

This enables one-click conversion of assignments to presentations!
"""

from typing import List, Dict, Any, Tuple
from dataclasses import dataclass
import re


@dataclass
class Slide:
    """Represents a single slide"""
    slide_number: int
    title: str
    layout: str  # "title", "bullet", "content", "conclusion"
    content: List[str] | str  # Content varies by layout
    speaker_notes: str = ""


def extract_sections(assignment_text: str) -> Dict[str, str]:
    """
    Extract markdown sections from assignment text.
    
    Looks for ## Section Header patterns in markdown.
    
    Args:
        assignment_text: Full assignment markdown text
        
    Returns:
        Dictionary mapping section names to their content
        
    Example:
        text = '''
        ## Introduction
        This is intro text.
        
        ## Methodology
        This is methodology.
        '''
        
        extract_sections(text) →
        {
            "Introduction": "This is intro text.\n",
            "Methodology": "This is methodology.\n"
        }
    """
    sections = {}
    current_section = None
    lines = assignment_text.split("\n")

    for line in lines:
        # Match ## Section Header pattern
        if line.startswith("## "):
            current_section = line.replace("## ", "").strip()
            sections[current_section] = ""
        elif current_section is not None:
            sections[current_section] += line + "\n"

    # Clean up sections
    for section_name in sections:
        sections[section_name] = sections[section_name].strip()

    return sections


def extract_key_points(text: str, max_points: int = 5) -> List[str]:
    """
    Extract key points from section text.
    
    Converts sentences into bullet points.
    
    Args:
        text: Section content text
        max_points: Maximum number of bullets to extract
        
    Returns:
        List of bullet point strings (max_points items)
        
    Example:
        text = "This is point one. This is point two. This is point three."
        extract_key_points(text, 2) → 
        ["This is point one.", "This is point two."]
    """
    # Split by sentences (. ! ?)
    sentences = re.split(r'[\.\!\?]+', text)
    
    # Clean up sentences
    bullets = []
    for sentence in sentences:
        sentence = sentence.strip()
        if len(sentence) > 10:  # Skip very short sentences
            # Capitalize and clean
            sentence = sentence[0].upper() + sentence[1:] if sentence else ""
            # Remove existing bullet points if present
            sentence = re.sub(r'^[\-\•\*]\s*', '', sentence)
            bullets.append(sentence)
    
    # Return up to max_points bullets
    return bullets[:max_points]


def create_slides_from_sections(topic: str, sections: Dict[str, str]) -> List[Dict[str, Any]]:
    """
    Convert assignment sections into slide structure.
    
    Creates a complete slide deck from assignment sections.
    
    Args:
        topic: Presentation topic
        sections: Dictionary of section name → content (from extract_sections)
        
    Returns:
        List of slide dictionaries ready for PPT export
        
    Output Format:
        [
            {
                "slide_number": 1,
                "title": "Topic",
                "layout": "title",
                "content": "Presentation Overview",
                "speaker_notes": "..."
            },
            {
                "slide_number": 2,
                "title": "Introduction",
                "layout": "bullet",
                "content": ["Point 1", "Point 2", ...],
                "speaker_notes": "..."
            },
            ...
        ]
    """
    slides = []
    slide_number = 1

    # Slide 1: Title Slide
    slides.append({
        "slide_number": slide_number,
        "title": topic,
        "layout": "title",
        "content": "Comprehensive Analysis & Presentation",
        "subtitle": "Research-Based Overview",
        "speaker_notes": f"Welcome to the presentation on {topic}. "
                        f"This presentation synthesizes key findings and research insights."
    })
    slide_number += 1

    # Slide 2: Overview/Agenda
    section_names = list(sections.keys())
    slides.append({
        "slide_number": slide_number,
        "title": "Agenda",
        "layout": "bullet",
        "content": section_names[:6],  # Show up to 6 sections in agenda
        "speaker_notes": f"We will cover {len(section_names)} main sections in this presentation."
    })
    slide_number += 1

    # Slides for each section
    for section_name, section_content in sections.items():
        if section_name.lower() in ["references", "title", "abstract"]:
            # Skip references, handle differently
            continue

        # Extract bullet points from section
        bullets = extract_key_points(section_content, max_points=5)

        if bullets:
            slides.append({
                "slide_number": slide_number,
                "title": section_name,
                "layout": "bullet",
                "content": bullets,
                "speaker_notes": f"Key points about {section_name}: "
                                + ". ".join(bullets[:2]) + "."
            })
            slide_number += 1

    # Final Slide: Conclusion
    slides.append({
        "slide_number": slide_number,
        "title": "Key Takeaways",
        "layout": "conclusion",
        "content": [
            "Research-backed insights",
            "Practical implications",
            "Future direction and recommendations",
            "Questions & Discussion"
        ],
        "speaker_notes": "Thank you for your attention. "
                        "The research presented here demonstrates the current state "
                        "of knowledge and highlights opportunities for future development."
    })

    return slides


def slides_to_json(slides: List[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Convert slides list to structured JSON format.
    
    Args:
        slides: List of slide dictionaries
        
    Returns:
        JSON-serializable dictionary
    """
    return {
        "slides": slides,
        "total_slides": len(slides),
        "presentation_format": "json",
        "slide_format": {
            "title": "Slide title",
            "layout": "bullet | title | content | conclusion",
            "content": "Content varies by layout",
            "speaker_notes": "Notes for presenter"
        }
    }


def build_presentation(
    topic: str,
    assignment_text: str
) -> Dict[str, Any]:
    """
    Main function: Convert assignment to presentation slides.
    
    Complete pipeline:
    1. Extract sections from assignment markdown
    2. Create slide structure from sections
    3. Format as JSON for export
    
    Args:
        topic: Presentation topic
        assignment_text: Full assignment markdown text
        
    Returns:
        Dictionary with:
        - slides: List of slide data
        - total_slides: Number of slides created
        - presentation_format: Format string ("json")
        - metadata: Additional info
        
    Example:
        assignment = db.query(Assignment).first()
        
        presentation = build_presentation(
            topic="AI in Healthcare",
            assignment_text=assignment.content
        )
        
        slides = presentation["slides"]  # List of slides
        total = presentation["total_slides"]  # Number of slides
    """
    # Step 1: Extract sections from markdown
    sections = extract_sections(assignment_text)

    if not sections:
        # Fallback if no sections found
        return {
            "slides": [],
            "total_slides": 0,
            "error": "No sections found in assignment text",
            "presentation_format": "json"
        }

    # Step 2: Create slide structure
    slides = create_slides_from_sections(topic, sections)

    # Step 3: Format as JSON
    presentation_data = slides_to_json(slides)

    # Add metadata
    presentation_data["metadata"] = {
        "topic": topic,
        "sections_found": len(sections),
        "section_names": list(sections.keys()),
        "generation_status": "success",
        "ready_for_export": True
    }

    return presentation_data


# ============ HELPER FUNCTIONS FOR PPT EXPORT ============

def slides_to_python_pptx(slides: List[Dict[str, Any]], topic: str) -> Any:
    """
    Convert slide data to python-pptx Presentation object.
    
    This function integrates with python-pptx library to create actual PowerPoint files.
    
    Args:
        slides: List of slide dictionaries
        topic: Presentation title
        
    Returns:
        python-pptx Presentation object (can be saved as .pptx)
        
    Note:
        Requires: pip install python-pptx
        
    Example:
        slides_list = [...]
        prs = slides_to_python_pptx(slides_list, "My Topic")
        prs.save("presentation.pptx")
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = slide_data.get("title", "")

        # Add content based on layout
        layout_type = slide_data.get("layout", "bullet")

        if layout_type == "title":
            # Title slide layout
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = slide_data.get("content", "")

        elif layout_type == "bullet":
            # Bullet point layout
            content = slide_data.get("content", [])
            if len(slide.placeholders) > 1:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()

                for bullet_text in content[:5]:  # Limit to 5 bullets
                    p = text_frame.add_paragraph()
                    p.text = str(bullet_text)
                    p.level = 0

        elif layout_type == "conclusion":
            # Conclusion layout with large text
            if len(slide.placeholders) > 1:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()

                for point in slide_data.get("content", [])[:4]:
                    p = text_frame.add_paragraph()
                    p.text = str(point)
                    p.level = 0
                    p.font.size = Pt(24)

    return prs


def save_presentation(prs: Any, filename: str) -> bool:
    """
    Save presentation to file.
    
    Args:
        prs: python-pptx Presentation object
        filename: Output file path (should end with .pptx)
        
    Returns:
        True if successful, False otherwise
    """
    try:
        prs.save(filename)
        return True
    except Exception as e:
        print(f"Error saving presentation: {e}")
        return False


# ============ CONVENIENCE FUNCTIONS ============

def build_and_export_presentation(
    topic: str,
    assignment_text: str,
    output_file: str
) -> Tuple[bool, str]:
    """
    Complete pipeline: Build presentation from assignment and export to PPT.
    
    Args:
        topic: Presentation topic
        assignment_text: Assignment markdown content
        output_file: Output file path (.pptx)
        
    Returns:
        Tuple of (success: bool, message: str)
        
    Example:
        success, message = build_and_export_presentation(
            topic="AI in Healthcare",
            assignment_text=assignment.content,
            output_file="presentation.pptx"
        )
        
        if success:
            print(f"✅ {message}")
        else:
            print(f"❌ {message}")
    """
    try:
        # Step 1: Build presentation structure
        presentation_data = build_presentation(topic, assignment_text)

        if presentation_data.get("error"):
            return False, f"Error building presentation: {presentation_data['error']}"

        slides = presentation_data.get("slides", [])

        if not slides:
            return False, "No slides generated from assignment"

        # Step 2: Convert to PowerPoint
        prs = slides_to_python_pptx(slides, topic)

        # Step 3: Save to file
        success = save_presentation(prs, output_file)

        if success:
            total_slides = len(slides)
            return True, f"✅ Presentation created: {output_file} ({total_slides} slides)"
        else:
            return False, "Failed to save presentation file"

    except Exception as e:
        return False, f"Error: {str(e)}"
